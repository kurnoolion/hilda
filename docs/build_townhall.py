"""
HILDA Townhall deck generator — 3 slides, casual lab style.
Run: python docs/build_townhall.py
Output: docs/HILDA_Townhall.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x2C, 0x3E, 0x50)   # dark charcoal-navy  — header fills
SLATE     = RGBColor(0x54, 0x6E, 0x7A)   # blue-grey slate     — accents
MID       = RGBColor(0x78, 0x90, 0x9C)   # medium blue-grey    — secondary text
LIGHT_BG  = RGBColor(0xEC, 0xEF, 0xF1)   # near-white grey     — card fills
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT      = RGBColor(0x37, 0x47, 0x4F)   # dark slate          — body text

# ── Helpers ────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def add_text(slide, l, t, w, h, text,
             size=12, bold=False, italic=False,
             color=TEXT, align=PP_ALIGN.LEFT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.15, NAVY)
    add_text(slide, 0.4, 0.12, 12.5, 0.7, title,
             size=28, bold=True, color=WHITE, font='Calibri Light')
    add_rect(slide, 0, 1.15, 13.33, 0.06, SLATE)
    if subtitle:
        add_text(slide, 0.4, 1.27, 12.5, 0.38, subtitle,
                 size=13, italic=True, color=SLATE)

def card(slide, x, y, w, h, title, body,
         title_size=12, body_size=10.5):
    add_rect(slide, x, y, w, h, LIGHT_BG)
    add_rect(slide, x, y, 0.07, h, SLATE)
    add_text(slide, x+0.15, y+0.07, w-0.22, 0.28,
             title, size=title_size, bold=True, color=NAVY)
    add_text(slide, x+0.15, y+0.35, w-0.22, h-0.4,
             body, size=body_size, color=TEXT)

# ── Presentation ───────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — The Problem
# ══════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
header(s1, "How We Work Today",
       "A look at the day-to-day reality for PMs managing connected-device programs")

pain = [
    ("Duplicated effort",
     "Every new device × customer combination starts from scratch — a fresh Excel built manually or copied and adapted from a previous one."),
    ("No visibility",
     "Status is stale by the time it's shared. Stakeholders email PMs for updates that should be self-serve in a live dashboard."),
    ("Onboarding takes too long",
     "A new PM or new device starts from zero. The process and context live in one person's head — not in the system."),
    ("Context scattered across channels",
     "The same delivery item has threads in email, messenger, and an issue tracker. PMs mentally stitch it all together every day."),
    ("Customer submission is fully manual",
     "PMs log into each customer system individually, format files per their requirements, and upload one by one."),
]

# 3 cards left column, 2 cards right column
CW, CH = 5.9, 0.97
LX, RX, Y0, GAP = 0.35, 6.9, 1.75, 0.11

for i, (t, b) in enumerate(pain):
    x = LX if i < 3 else RX
    y = Y0 + (i if i < 3 else i - 3) * (CH + GAP)
    card(s1, x, y, CW, CH, t, b)

# Summary bar
add_rect(s1, 0, 7.06, 13.33, 0.44, LIGHT_BG)
add_text(s1, 0.4, 7.10, 12.5, 0.35,
         "Result: high PM toil  ·  slow turnaround  ·  no scalability to new customers or new PMs without growing headcount",
         size=11, italic=True, color=MID, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What HILDA Does
# ══════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
header(s2, "Introducing HILDA / DeliverableHub",
       "Template-driven automation for the full deliverable lifecycle — PM manages by exception, not by spreadsheet")

# Pipeline boxes
steps = [
    ("1  Template",  "PM selects customer\ntemplate; full tracker\nauto-generated"),
    ("2  Kickoff",   "One click sends\nstructured requests\nto every owner"),
    ("3  Track",     "Automation captures\nresponses & updates\ndashboard live"),
    ("4  Review",    "AI first-pass quality\ncheck; PM resolves\nopen items"),
    ("5  Submit",    "Package assembled\nper customer format;\nPM previews & approves"),
]

BW, BH, BY = 2.15, 1.65, 1.85
ARROW = 0.22
total = len(steps) * BW + (len(steps)-1) * ARROW
sx = (13.33 - total) / 2

for i, (st, sd) in enumerate(steps):
    bx = sx + i * (BW + ARROW)
    add_rect(s2, bx, BY, BW, BH, NAVY)
    add_text(s2, bx, BY+0.1, BW, 0.35, st,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s2, bx+0.2, BY+0.47, BW-0.4, 0.03, SLATE)   # divider line
    add_text(s2, bx+0.1, BY+0.55, BW-0.2, 1.0, sd,
             size=10, color=LIGHT_BG, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        add_text(s2, bx+BW, BY+0.65, ARROW, 0.35, "›",
                 size=18, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

# Human-in-the-loop callout
add_rect(s2, 0.7, 3.72, 11.93, 0.68, LIGHT_BG)
add_rect(s2, 0.7, 3.72, 0.07, 0.68, SLATE)
add_text(s2, 0.85, 3.78, 11.7, 0.55,
         "Human-in-the-loop  —  Every customer-facing action requires explicit PM sign-off."
         "  HILDA assists; it never acts on your behalf without approval.",
         size=12, color=NAVY, bold=False)

# Role cards
roles = [
    ("PMs / TPMs",
     "Dashboard replaces Excel. Approvals replace copy-paste. All channels in one place."),
    ("R&D owners",
     "Structured requests with clear reference tags. Fewer follow-up emails chasing status."),
    ("PM team leads",
     "Author templates once. Every device and new PM inherits them — no bespoke trackers."),
]

add_text(s2, 0.35, 4.55, 5, 0.3, "What changes for you",
         size=13, bold=True, color=NAVY)

RCW, RCH = 4.05, 0.88
for i, (rt, rb) in enumerate(roles):
    rx = 0.35 + i * (RCW + 0.32)
    card(s2, rx, 4.88, RCW, RCH, rt, rb, title_size=12, body_size=10)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Status & Roadmap
# ══════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
header(s3, "Current Status & What's Coming",
       "Dates are tentative — we will keep the team updated as targets firm up")

COL_W = 5.95
C1X, C2X = 0.35, 6.73

# Ph-1 header
add_rect(s3, C1X, 1.72, COL_W, 0.52, NAVY)
add_text(s3, C1X, 1.72, COL_W, 0.52,
         "Phase 1  ·  Target: 30 Jun 2026 (tentative)",
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

ph1 = [
    "SharePoint dashboard — live delivery status per device & milestone",
    "Email ingest — structured + free-text reply routing to delivery items",
    "Shared network drive — owner document drop & auto-classification",
    "Corp PLM polling — document download & indexing",
    "Rule-based reminders & escalation (configurable schedule)",
    "Jira (internal tracker) + Slack messenger adapters",
    "Test report parser + final / interim classifier",
    "AI quality review — first-pass findings on reports & waivers",
]

y = 2.33
for item in ph1:
    add_text(s3, C1X+0.15, y, COL_W-0.2, 0.31,
             f"•  {item}", size=10.5, color=TEXT)
    y += 0.32

# Ph-2 header
add_rect(s3, C2X, 1.72, COL_W, 0.52, SLATE)
add_text(s3, C2X, 1.72, COL_W, 0.52,
         "Phase 2  ·  Target: 31 Jul 2026 (tentative)",
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

ph2 = [
    "Revision tracking — multi-version document lifecycle per item",
    "Version selection workflow via corp messenger",
    "PLM sync & stale attachment cleanup post-submission",
    "Customer portal submission — assembled package, PM sign-off",
    "Proprietary internal messenger adapter",
    "PM resolution-path annotation for failed test cases",
    "SharePoint document upload surface (PM / TPM)",
    "Per-item owner self-close from SharePoint UI",
]

y = 2.33
for item in ph2:
    add_text(s3, C2X+0.15, y, COL_W-0.2, 0.31,
             f"◦  {item}", size=10.5, color=TEXT)
    y += 0.32

# Footer
add_rect(s3, 0, 7.06, 13.33, 0.44, LIGHT_BG)
add_text(s3, 0.4, 7.10, 12.5, 0.35,
         "Questions or feedback? Reach out to the HILDA team  "
         "·  Detailed walkthroughs coming as we approach Phase 1 target",
         size=11, italic=True, color=MID, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────
out = r'C:\projects\hilda\hilda\docs\HILDA_Townhall.pptx'
prs.save(out)
print(f"Saved: {out}")
